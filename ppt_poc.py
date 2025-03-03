import streamlit as st
import pandas as pd
import numpy as np
import altair as alt
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO

# Enhanced Sales Data with more attributes
def create_enhanced_data():
    np.random.seed(42)  # For reproducibility
    
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
    regions = ['North', 'South', 'East', 'West']
    products = ['Product A', 'Product B', 'Product C']
    
    data = []
    
    for month in months:
        for region in regions:
            for product in products:
                # Base sales
                sales = np.random.randint(15000, 50000)
                # Cost (60-80% of sales)
                cost = int(sales * np.random.uniform(0.6, 0.8))
                # Profit
                profit = sales - cost
                # Margin percentage
                margin = round((profit / sales) * 100, 1)
                # Units sold
                unit_price = np.random.uniform(50, 200)
                units = int(sales / unit_price)
                # Market share
                market_share = np.random.uniform(5, 25)
                # Customer satisfaction
                csat = np.random.randint(70, 98)
                
                data.append({
                    'Month': month,
                    'Region': region,
                    'Product': product,
                    'Sales': sales,
                    'Cost': cost,
                    'Profit': profit,
                    'Margin': margin,
                    'Units': units,
                    'UnitPrice': round(unit_price, 2),
                    'MarketShare': round(market_share, 1),
                    'CSAT': csat
                })
    
    return pd.DataFrame(data)

# Function to create aggregated datasets for different chart types
def prepare_chart_datasets(df, dimension, measure):
    """
    Aggregate data based on selected dimension and measure
    
    Parameters:
    - df: Full dataset
    - dimension: Column to group by (e.g., 'Month', 'Region', 'Product')
    - measure: Column to aggregate (e.g., 'Sales', 'Profit', 'Units')
    
    Returns:
    - Aggregated DataFrame
    """
    # Group by the selected dimension
    aggregated = df.groupby(dimension)[measure].sum().reset_index()
    
    # If dimension is Month, sort it chronologically
    if dimension == 'Month':
        month_order = {'Jan': 1, 'Feb': 2, 'Mar': 3, 'Apr': 4, 'May': 5, 'Jun': 6}
        aggregated['MonthNum'] = aggregated['Month'].map(month_order)
        aggregated = aggregated.sort_values('MonthNum')
        if 'MonthNum' in aggregated.columns:
            aggregated = aggregated.drop('MonthNum', axis=1)
    
    return aggregated

# Function to create PowerPoint with dynamic chart types
def create_ppt(df, chart_type, dimension, measure):
    # Create a presentation object
    prs = Presentation()
    
    # Make sure we're working with clean data
    df = df.reset_index(drop=True).copy()
    
    # Create aggregated dataset based on selected dimension and measure
    chart_data_df = prepare_chart_datasets(df, dimension, measure)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{measure} by {dimension}"
    subtitle.text = f"Generated from Streamlit App - March 2025"
    
    # Dataset Slide - separate slide for the data
    dataset_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(dataset_slide_layout)
    title = slide.shapes.title
    title.text = f"{measure} by {dimension} Data"
    
    # Add table
    rows, cols = len(chart_data_df) + 1, len(chart_data_df.columns)
    left = Inches(1.5)
    top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(4.0)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    for col_idx, col_name in enumerate(chart_data_df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light gray background
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)  # Standard size for PowerPoint
    
    # Data
    for row_idx, row in chart_data_df.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            # Make measure values blue to indicate they're editable
            if col_idx == 1:  # Measure column (second column)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor(0, 0, 220)  # Blue color
                paragraph.font.bold = True
    
    # Chart Slide - separate slide for the chart
    chart_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(chart_slide_layout)
    title = slide.shapes.title
    title.text = f"{measure} by {dimension}"
    
    # Create chart with the selected data
    chart_data = CategoryChartData()
    
    # Add categories and series data
    chart_data.categories = chart_data_df[dimension].tolist()
    chart_data.add_series(measure, chart_data_df[measure].tolist())
    
    # Adjust chart dimensions based on chart type
    if chart_type == XL_CHART_TYPE.PIE:
        x, y, cx, cy = Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.5)
    # Special handling for line charts to avoid repair prompts
    elif chart_type in [XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS]:
        # For line charts, use standard dimensions and different positioning
        x, y, cx, cy = Inches(1.0), Inches(1.5), Inches(6.0), Inches(4.0)
    else:
        # For other charts like column and bar
        x, y, cx, cy = Inches(0.75), Inches(1.5), Inches(6.5), Inches(4.5)
    
    # Create the chart with appropriate type and dimensions
    chart = slide.shapes.add_chart(
        chart_type, x, y, cx, cy, chart_data
    ).chart
    
    # Configure chart with minimal formatting to avoid repair issues
    chart.has_title = True
    chart.chart_title.text_frame.text = f"{measure} by {dimension}"
    
    # Always place legend outside the chart
    chart.has_legend = True
    if chart_type == XL_CHART_TYPE.PIE:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    else:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
    
    # Do not use include_in_layout for line charts as it can cause repair issues
    if chart_type not in [XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS]:
        chart.legend.include_in_layout = False
    
    # Axis formatting for non-pie charts only
    if chart_type != XL_CHART_TYPE.PIE:
        try:
            category_axis = chart.category_axis
            if hasattr(category_axis, 'tick_labels'):
                category_axis.tick_labels.font.size = Pt(12)  # Safe size
                category_axis.tick_labels.font.bold = True
        except (AttributeError, ValueError):
            # Skip axis formatting if not applicable
            pass
    
    # Data labels - special handling for line charts
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(12)  # Safe size
    data_labels.font.bold = True
    data_labels.number_format = '#,##0'
    
    # Different data label position for different chart types
    if chart_type == XL_CHART_TYPE.PIE:
        # For pie charts, let PowerPoint handle label positioning
        pass
    elif chart_type in [XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS]:
        # For line charts, above position works better and causes fewer repair issues
        data_labels.position = XL_LABEL_POSITION.ABOVE
    else:
        # For other charts like column and bar
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    
    # Add instructions for updating
    left = Inches(1.0)
    top = Inches(6.5)
    width = Inches(7.0)
    height = Inches(0.5)
    txtbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txtbox.text_frame
    p = tf.add_paragraph()
    p.text = "Right-click chart and select 'Edit Data' to update after changing values in the data table"
    p.font.italic = True
    p.font.size = Pt(12)
    
    # Save to BytesIO buffer
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)  # Reset buffer position
    return ppt_buffer

# Map friendly names to PowerPoint chart types
CHART_TYPES = {
    "Column Chart": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "Bar Chart": XL_CHART_TYPE.BAR_CLUSTERED,
    "Line Chart": XL_CHART_TYPE.LINE,
    "Line Chart with Markers": XL_CHART_TYPE.LINE_MARKERS,
    "Pie Chart": XL_CHART_TYPE.PIE,
    "Area Chart": XL_CHART_TYPE.AREA
}

# Streamlit App
def main():
    st.title("Enhanced Sales Dashboard to PowerPoint Converter")
    
    # Generate and cache enhanced data
    @st.cache_data
    def get_data():
        return create_enhanced_data()
    
    df = get_data()
    
    # Dashboard configuration
    st.sidebar.header("Dashboard Configuration")
    
    # Chart type selection
    chart_type_name = st.sidebar.selectbox(
        "Select Chart Type", 
        list(CHART_TYPES.keys())
    )
    chart_type = CHART_TYPES[chart_type_name]
    
    # Dimension selection (what to group by)
    dimension_options = ['Month', 'Region', 'Product']
    dimension = st.sidebar.selectbox(
        "Group By (Dimension)", 
        dimension_options
    )
    
    # Measure selection (what to calculate)
    measure_options = ['Sales', 'Cost', 'Profit', 'Units', 'Margin', 'MarketShare', 'CSAT']
    measure = st.sidebar.selectbox(
        "Measure", 
        measure_options
    )
    
    # Show data preview with filters
    with st.expander("View Data", expanded=False):
        st.dataframe(df, use_container_width=True)
    
    # Generate aggregated data for selected dimension and measure
    chart_data_df = prepare_chart_datasets(df, dimension, measure)
    
    # Show aggregated data
    st.subheader(f"{measure} by {dimension}")
    st.dataframe(chart_data_df, use_container_width=True)
    
    # Display chart in Streamlit with data labels using Altair
    st.subheader("Chart Preview")
    
    # Create appropriate Altair chart based on selection
    if chart_type_name in ["Column Chart"]:
        chart = alt.Chart(chart_data_df).mark_bar().encode(
            x=alt.X(f'{dimension}:N', title=dimension),
            y=alt.Y(f'{measure}:Q', title=measure)
        )
    elif chart_type_name in ["Bar Chart"]:
        chart = alt.Chart(chart_data_df).mark_bar().encode(
            y=alt.Y(f'{dimension}:N', title=dimension),
            x=alt.X(f'{measure}:Q', title=measure)
        )
    elif chart_type_name in ["Line Chart", "Line Chart with Markers"]:
        chart = alt.Chart(chart_data_df).mark_line(
            point=chart_type_name == "Line Chart with Markers"
        ).encode(
            x=alt.X(f'{dimension}:N', title=dimension),
            y=alt.Y(f'{measure}:Q', title=measure)
        )
    elif chart_type_name == "Pie Chart":
        # For pie charts, we need a different approach
        chart = alt.Chart(chart_data_df).mark_arc().encode(
            theta=alt.Theta(f'{measure}:Q'),
            color=alt.Color(f'{dimension}:N')
        )
    elif chart_type_name == "Area Chart":
        chart = alt.Chart(chart_data_df).mark_area().encode(
            x=alt.X(f'{dimension}:N', title=dimension),
            y=alt.Y(f'{measure}:Q', title=measure)
        )
    else:
        # Default to column chart
        chart = alt.Chart(chart_data_df).mark_bar().encode(
            x=alt.X(f'{dimension}:N', title=dimension),
            y=alt.Y(f'{measure}:Q', title=measure)
        )
    
    # Add text labels for all non-pie charts
    if chart_type_name != "Pie Chart":
        text = chart.mark_text(
            align='center',
            baseline='bottom',
            dy=-5,
            fontSize=14,
            color='black'
        ).encode(
            text=alt.Text(f'{measure}:Q', format=',.0f')
        )
        chart = chart + text
    
    # Adjust chart properties to match PowerPoint layout
    if chart_type_name == "Pie Chart":
        chart_properties = {
            'width': 650,
            'height': 400,
            'title': f'{measure} by {dimension}'
        }
    else:
        # Make the chart narrower to match PowerPoint with legend on right
        chart_properties = {
            'width': 550,  # Reduced width to match PowerPoint proportions with legend
            'height': 400,
            'title': f'{measure} by {dimension}'
        }
    
    # Configure legend position
    if chart_type_name == "Pie Chart":
        final_chart = chart.properties(**chart_properties).configure_legend(
            orient='bottom'
        )
    else:
        final_chart = chart.properties(**chart_properties).configure_legend(
            orient='right',
            offset=20
        )
    
    st.altair_chart(final_chart, use_container_width=True)
    
    # Download button for PPT
    if st.button("Generate PowerPoint"):
        with st.spinner("Creating PowerPoint..."):
            ppt_buffer = create_ppt(df, chart_type, dimension, measure)
        
        st.success("PowerPoint generated successfully!")
        
        st.download_button(
            label="Download PowerPoint",
            data=ppt_buffer.getvalue(),
            file_name=f"{measure}_by_{dimension}_{chart_type_name}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
        # Usage instructions
        st.info("""
        **To use the PowerPoint:**
        1. Open the downloaded file
        2. View the data table on the second slide
        3. Edit the blue values if needed
        4. Go to the chart slide and right-click the chart
        5. Select 'Edit Data' to update based on your changes
        """)

if __name__ == "__main__":
    main()