from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
import plotly.graph_objects as go
from io import BytesIO
from PIL import Image

# Create a new PowerPoint Presentation
prs = Presentation()

### Function to Add Styled Headings ###
def add_heading(slide, text, font_size=32):
    title_shape = slide.shapes.title
    title_shape.text = text
    title_shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue

### SLIDE 1: Title Slide ###
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Automated PowerPoint Report"
subtitle.text = "Generated using Python (python-pptx + Plotly)"

### SLIDE 2: Narrative Slide ###
slide_layout = prs.slide_layouts[1]  # Title + Content
slide = prs.slides.add_slide(slide_layout)
add_heading(slide, "Dynamic Narrative")
content = slide.placeholders[1]
content.text = (
    "This slide demonstrates how Python can generate narratives.\n\n"
    "Using `python-pptx`, we can automate reports, dashboards, and data presentations dynamically."
)

### SLIDE 3: Table with Header Highlight ###
slide_layout = prs.slide_layouts[5]  # Title + Table
slide = prs.slides.add_slide(slide_layout)
add_heading(slide, "Data Table with Styling")

# Sample Data
df = pd.DataFrame({
    "Category": ["A", "B", "C", "D"],
    "Value": [100, 250, 180, 300]
})

# Create Table
rows, cols = df.shape
table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(2), Inches(6), Inches(2)).table

# Add Header Row with Styling
for col_idx, col_name in enumerate(df.columns):
    cell = table.cell(0, col_idx)
    cell.text = col_name
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark Blue
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    cell.text_frame.paragraphs[0].font.bold = True

# Add Table Data
for row_idx, row in df.iterrows():
    for col_idx, value in enumerate(row):
        table.cell(row_idx + 1, col_idx).text = str(value)

### SLIDE 4: Plotly Chart Slide ###
slide_layout = prs.slide_layouts[5]  # Title + Content
slide = prs.slides.add_slide(slide_layout)
add_heading(slide, "Interactive Chart (Plotly)")

# Create a Plotly Bar Chart
fig = go.Figure(data=[
    go.Bar(x=df["Category"], y=df["Value"], marker_color="blue")
])
fig.update_layout(title="Category Values", xaxis_title="Category", yaxis_title="Value")

# Save the chart as an image
img_stream = BytesIO()
fig.write_image(img_stream, format="png")
img_stream.seek(0)

# Insert image into PPT
img = slide.shapes.add_picture(img_stream, Inches(2), Inches(2), Inches(6), Inches(4))

### SLIDE 5: Pie Chart with Plotly ###
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
add_heading(slide, "Pie Chart (Plotly)")

# Generate a Pie Chart
fig_pie = go.Figure(data=[go.Pie(labels=df["Category"], values=df["Value"], hole=0.3)])
fig_pie.update_layout(title="Category Distribution")

# Save the pie chart as an image
img_pie_stream = BytesIO()
fig_pie.write_image(img_pie_stream, format="png")
img_pie_stream.seek(0)

# Insert pie chart image
slide.shapes.add_picture(img_pie_stream, Inches(2), Inches(2), Inches(6), Inches(4))

### SLIDE 6: Custom Shape ###
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
add_heading(slide, "Shapes and Styling")

# Add a Rectangle Shape
shape = slide.shapes.add_shape(
    autoshape_type=1,  # Rectangle shape
    left=Inches(2),
    top=Inches(2),
    width=Inches(4),
    height=Inches(1)
)
shape.text = "Python-Powered Report"
shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White

### SAVE PRESENTATION ###
prs.save("Enhanced_Presentation.pptx")
print("Presentation successfully created with advanced visuals!")
